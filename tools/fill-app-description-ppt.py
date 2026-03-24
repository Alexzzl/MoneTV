from __future__ import annotations

from pathlib import Path
from typing import Iterable

from PIL import Image, ImageDraw, ImageFont, ImageOps
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path.cwd()
TEMPLATE = ROOT / "App_Description_template_eng_v.1.42.pptx"
APP_NAME = "MoneTV"
COMPANY_NAME = "MATICONE LTD"
OUTPUT = ROOT / f"App_Description_{APP_NAME}_filled.pptx"
SCREENSHOT_DIR = ROOT / "store-screenshots"
UI_DIR = ROOT / "ui-description-screenshots"


def remove_slide(prs: Presentation, index: int) -> None:
    slide_id = prs.slides._sldIdLst[index]
    prs.part.drop_rel(slide_id.rId)
    prs.slides._sldIdLst.remove(slide_id)


def remove_shape(shape) -> None:
    shape._element.getparent().remove(shape._element)


def clear_and_set_text(shape, lines: Iterable[str], font_size: int = 20, bold: bool = False) -> None:
    text_frame = shape.text_frame
    text_frame.clear()

    for idx, line in enumerate(lines):
        paragraph = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
        paragraph.text = line
        for run in paragraph.runs:
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = RGBColor(0x22, 0x22, 0x22)


def set_cell_text(cell, text: str, font_size: int = 12, bold: bool = False) -> None:
    cell.text = text
    for paragraph in cell.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(font_size)
            run.font.bold = bold


def find_font(size: int) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    candidates = [
        Path("C:/Windows/Fonts/arial.ttf"),
        Path("C:/Windows/Fonts/calibri.ttf"),
    ]
    for candidate in candidates:
        if candidate.exists():
            return ImageFont.truetype(str(candidate), size=size)
    return ImageFont.load_default()


def create_collage(items: list[tuple[Path, str]], output_path: Path, title: str) -> None:
    width, height = 1600, 900
    margin = 40
    gutter = 24
    title_height = 70
    cols, rows = 2, 2
    cell_width = (width - margin * 2 - gutter) // cols
    cell_height = (height - margin * 2 - title_height - gutter) // rows
    image_height = cell_height - 52

    canvas = Image.new("RGB", (width, height), (248, 248, 248))
    draw = ImageDraw.Draw(canvas)
    title_font = find_font(34)
    label_font = find_font(24)

    draw.text((margin, 18), title, fill=(28, 28, 28), font=title_font)

    for idx, (image_path, label) in enumerate(items):
        row = idx // cols
        col = idx % cols
        x = margin + col * (cell_width + gutter)
        y = margin + title_height + row * (cell_height + gutter)

        image = Image.open(image_path).convert("RGB")
        fitted = ImageOps.fit(image, (cell_width, image_height), method=Image.Resampling.LANCZOS)
        canvas.paste(fitted, (x, y))

        label_top = y + image_height
        draw.rectangle((x, label_top, x + cell_width, y + cell_height), fill=(255, 255, 255))
        draw.rectangle((x, y, x + cell_width, y + cell_height), outline=(210, 210, 210), width=2)
        draw.text((x + 16, label_top + 12), label, fill=(40, 40, 40), font=label_font)

    canvas.save(output_path, quality=90)


def add_centered_box(slide, left, top, width, height, text: str, fill_rgb: tuple[int, int, int]) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*fill_rgb)
    shape.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    text_frame = shape.text_frame
    text_frame.clear()
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


def add_arrow_text(slide, left, top, width, height, text: str) -> None:
    shape = slide.shapes.add_textbox(left, top, width, height)
    text_frame = shape.text_frame
    text_frame.clear()
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x55, 0x55, 0x55)


def main() -> None:
    if not TEMPLATE.exists():
        raise FileNotFoundError(f"Template not found: {TEMPLATE}")

    popup_home = UI_DIR / "07-home-return-popup.jpg"
    usage_collage = UI_DIR / "08-usage-scenario-collage.jpg"
    menu_collage = UI_DIR / "09-menu-function-collage.jpg"
    return_collage = UI_DIR / "10-return-key-policy-collage.jpg"

    create_collage(
        [
            (SCREENSHOT_DIR / "01-home-hero.jpg", "1. Home screen"),
            (SCREENSHOT_DIR / "02-home-popular.jpg", "2. Browse drama cards"),
            (UI_DIR / "05-detail-page.jpg", "3. Detail screen"),
            (UI_DIR / "06-player-page.jpg", "4. Player screen"),
        ],
        usage_collage,
        f"{APP_NAME} - Usage Scenario",
    )

    create_collage(
        [
            (SCREENSHOT_DIR / "01-home-hero.jpg", "Home"),
            (SCREENSHOT_DIR / "03-discover-top.jpg", "Discover"),
            (UI_DIR / "05-detail-page.jpg", "Detail"),
            (UI_DIR / "06-player-page.jpg", "Player"),
        ],
        menu_collage,
        f"{APP_NAME} - Main Screens and Functions",
    )

    create_collage(
        [
            (SCREENSHOT_DIR / "01-home-hero.jpg", "1. Home screen before pressing Return"),
            (popup_home, "2. Press Return on Home: exit confirmation popup opens"),
            (SCREENSHOT_DIR / "03-discover-top.jpg", "3. Press Return on Discover / Detail / Player: previous page"),
            (SCREENSHOT_DIR / "01-home-hero.jpg", "4. Press Return again while popup is open: popup closes"),
        ],
        return_collage,
        f"{APP_NAME} - Return Key Policy",
    )

    prs = Presentation(str(TEMPLATE))
    remove_slide(prs, 0)

    cover = prs.slides[0]
    cover.shapes[0].text = f"{APP_NAME} Description"
    cover.shapes[1].text = COMPANY_NAME

    revision = prs.slides[1]
    revision.shapes[0].text = "Revision History"
    table = revision.shapes[1].table
    set_cell_text(table.cell(0, 0), "Version", 12, True)
    set_cell_text(table.cell(0, 1), "Date", 12, True)
    set_cell_text(table.cell(0, 2), "Description", 12, True)
    set_cell_text(table.cell(0, 3), "Author", 12, True)
    set_cell_text(table.cell(1, 0), "1.0")
    set_cell_text(table.cell(1, 1), "March 19, 2026")
    set_cell_text(table.cell(1, 2), f"Initial Samsung TV Seller Office UI description for {APP_NAME}.")
    set_cell_text(table.cell(1, 3), "db")
    set_cell_text(table.cell(2, 0), "1.0.1")
    set_cell_text(table.cell(2, 1), "March 24, 2026")
    set_cell_text(table.cell(2, 2), "Updated return key policy scenarios and screenshots for Samsung review.")
    set_cell_text(table.cell(2, 3), "db")

    contents = prs.slides[2]
    contents.shapes[0].text = "Contents"
    clear_and_set_text(
        contents.shapes[1],
        [
            "UI Structure",
            "Usage Scenario",
            "Menu & Function Description",
            "Key Policy",
            "How to Change Languages",
        ],
        font_size=24,
    )

    ui_overview = prs.slides[3]
    ui_overview.shapes[0].text = "UI Structure"
    clear_and_set_text(
        ui_overview.shapes[1],
        [
            "Whole UI Structure",
            "",
            f"{APP_NAME} is a Samsung TV video application for browsing short-drama content",
            "and playing selected episodes.",
            "Only currently available user flows are described in this document.",
            "",
            "Primary user flow:",
            "Home -> Detail -> Player",
            "Discover -> Detail -> Player",
            "",
            "Main screens in the current release:",
            "- Home: featured banner and content rows",
            "- Discover: trending, genre, and new-this-week sections",
            "- Detail: title information, episode list, and recommendation area",
            "- Player: video playback and episode list",
        ],
        font_size=18,
    )

    flow_slide = prs.slides[4]
    flow_slide.shapes[0].text = "UI Structure - Flow Graph"
    remove_shape(flow_slide.shapes[1])
    add_centered_box(flow_slide, Inches(1.9), Inches(1.9), Inches(1.7), Inches(0.65), "Home", (0x3F, 0xA7, 0x7A))
    add_arrow_text(flow_slide, Inches(3.75), Inches(1.95), Inches(0.5), Inches(0.5), "<->")
    add_centered_box(flow_slide, Inches(4.3), Inches(1.9), Inches(1.95), Inches(0.65), "Discover", (0x3F, 0xA7, 0x7A))
    add_arrow_text(flow_slide, Inches(3.95), Inches(2.75), Inches(0.3), Inches(0.5), "v")
    add_centered_box(flow_slide, Inches(3.25), Inches(3.2), Inches(1.8), Inches(0.7), "Detail", (0xE0, 0x8A, 0x3B))
    add_arrow_text(flow_slide, Inches(4.0), Inches(3.95), Inches(0.3), Inches(0.5), "v")
    add_centered_box(flow_slide, Inches(3.05), Inches(4.35), Inches(2.2), Inches(0.75), "Player", (0xB9, 0x4E, 0x5D))

    depth_slide = prs.slides[5]
    depth_slide.shapes[0].text = "UI Structure - Depth Navigation"
    remove_shape(depth_slide.shapes[1])
    textbox = depth_slide.shapes.add_textbox(Inches(1.15), Inches(2.0), Inches(7.1), Inches(4.1))
    clear_and_set_text(
        textbox,
        [
            "Depth 1 : Home, Discover",
            "Depth 2 : Selected drama detail page",
            "Depth 3 : Episode playback page",
            "",
            "Global UI elements:",
            "- Header navigation bar",
            "- Focus highlight for remote-control navigation",
            "",
            "Contextual UI elements:",
            "- Episode list on Detail page",
            "- Video area and episode list on Player page",
        ],
        font_size=20,
    )

    usage_summary = prs.slides[6]
    usage_summary.shapes[0].text = "Usage Scenario"
    clear_and_set_text(
        usage_summary.shapes[1],
        [
            "Use Case 1: Browse content and play an episode.",
            "",
            "1. Move focus on the Home or Discover screen with the directional keys.",
            "2. Select a drama card with the Enter key.",
            "3. Review the title information and episode list on the Detail screen.",
            "4. Select Play Episode 1 or choose an available episode card.",
            "5. The Player screen opens and video playback starts.",
            "6. Select another episode from the episode list if needed.",
            "7. Press Return on Discover / Detail / Player to go to the previous screen.",
            "",
            "No account sign-in, activation, or purchase is required in the current release.",
        ],
        font_size=18,
    )

    usage_sample = prs.slides[7]
    usage_sample.shapes[0].text = "Usage Scenario"
    clear_and_set_text(
        usage_sample.shapes[1],
        [
            "Use Case Title: Return key policy scenarios",
            "Path 1: Home -> Return -> Exit confirmation popup",
            "Path 2: Exit popup open -> Return -> Popup closes and Home remains active",
            "Path 3: Discover / Detail / Player -> Return -> Previous page",
            "",
            "Supporting notes:",
            "- The app follows Samsung TV Return key policy in the current release.",
            "- The Home page does not exit immediately on Return.",
            "- The dedicated EXIT key is left to the Samsung TV platform default behavior.",
            "",
            "No account login is required.",
            "No in-app purchase is required.",
        ],
        font_size=18,
    )
    usage_picture = usage_sample.shapes[2]
    left, top, width, height = usage_picture.left, usage_picture.top, usage_picture.width, usage_picture.height
    remove_shape(usage_picture)
    usage_sample.shapes.add_picture(str(return_collage), left, top, width=width, height=height)

    menu_slide = prs.slides[8]
    menu_slide.shapes[0].text = "Menu & Function Description"
    clear_and_set_text(
        menu_slide.shapes[1],
        [
            "Screen descriptions in the current release:",
            "",
            "Home: displays a featured banner and drama rows for browsing.",
            "Discover: displays trending content, genre browsing blocks, and new content sections.",
            "Detail: displays title metadata, synopsis, episode list, and recommendation area.",
            "Player: displays the selected video and the episode list for switching episodes.",
            "",
            "Only currently implemented screens and functions are described in this document.",
        ],
        font_size=18,
    )
    menu_picture = menu_slide.shapes[2]
    left, top, width, height = menu_picture.left, menu_picture.top, menu_picture.width, menu_picture.height
    remove_shape(menu_picture)
    menu_slide.shapes.add_picture(str(menu_collage), left, top, width=width, height=height)

    key_slide = prs.slides[9]
    key_slide.shapes[0].text = "Key Policy"
    clear_and_set_text(
        key_slide.shapes[1],
        [
            "The current release follows Samsung TV Return key policy.",
            "Home page: Return opens the exit confirmation popup.",
            "Exit popup open: Return closes the popup first.",
            "Discover / Detail / Player pages: Return navigates to the previous page.",
            "The dedicated EXIT key uses the Samsung TV platform default behavior.",
        ],
        font_size=18,
    )
    key_table = key_slide.shapes[2].table
    rows = [
        ("Button", "Action", "Remarks"),
        ("ENTER", "Select the focused card, button, or episode", "Standard selection behavior"),
        ("UP / DOWN", "Move focus vertically", "Standard directional navigation"),
        ("LEFT / RIGHT", "Move focus horizontally", "Standard directional navigation"),
        ("RETURN on Home", "Open the exit confirmation popup", "Required by Samsung TV Return key policy"),
        ("RETURN on popup", "Close the exit confirmation popup", "The app stays on the Home page"),
        ("RETURN on Discover / Detail / Player", "Go to the previous screen", "Navigates back within the app flow"),
        ("EXIT", "Close the application", "Platform default Samsung TV behavior"),
        ("PLAY / PAUSE / FF / RW", "N/R", "No custom key mapping in current release"),
        ("COLOR / CHANNEL / NUMBER KEYS", "N/R", "No custom key mapping in current release"),
    ]
    for r_idx, row in enumerate(rows):
        for c_idx, value in enumerate(row):
            set_cell_text(key_table.cell(r_idx, c_idx), value, 11, r_idx == 0)

    lang_slide = prs.slides[10]
    lang_slide.shapes[0].text = "How to Change Languages"
    lang_table = lang_slide.shapes[1].table
    set_cell_text(lang_table.cell(0, 0), "Items", 12, True)
    set_cell_text(lang_table.cell(0, 1), "Contents", 12, True)
    set_cell_text(lang_table.cell(1, 0), "Language support")
    set_cell_text(
        lang_table.cell(1, 1),
        "This version supports English only. No in-app language change menu is provided.",
    )

    prs.save(str(OUTPUT))
    print(f"saved={OUTPUT}")


if __name__ == "__main__":
    main()
